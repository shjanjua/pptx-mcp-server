"""
Apply text replacements to PowerPoint presentations.

This module provides functionality to:
- Replace text content in PowerPoint shapes using JSON specifications
- Preserve or modify paragraph formatting (alignment, bullets, fonts, spacing)
- Clear all text shapes unless explicitly provided with new content
- Validate replacements before applying them
"""

import json
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape
from pptx.util import Pt

from .inventory import (
    collect_shapes_with_absolute_positions,
    is_valid_shape,
    sort_shapes_by_position,
    ShapeData,
)


# Type aliases
JsonValue = Union[str, int, float, bool, None]
ParagraphSpec = Dict[str, JsonValue]
ReplacementSpec = Dict[str, Any]


def validate_shape_id(shape_id: str) -> bool:
    """Validate that a shape ID follows the expected format."""
    return bool(re.match(r"^shape-\d+$", shape_id))


def validate_slide_id(slide_id: str) -> bool:
    """Validate that a slide ID follows the expected format."""
    return bool(re.match(r"^slide-\d+$", slide_id))


def parse_slide_index(slide_id: str) -> int:
    """Extract slide index from slide ID string."""
    return int(slide_id.split("-")[1])


def parse_shape_index(shape_id: str) -> int:
    """Extract shape index from shape ID string."""
    return int(shape_id.split("-")[1])


def clear_text_frame(shape: BaseShape) -> None:
    """Clear all content from a shape's text frame."""
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return

    tf = shape.text_frame
    # Keep first paragraph but clear its content and remove extra paragraphs
    if tf.paragraphs:
        first_para = tf.paragraphs[0]
        first_para.clear()

        # Remove all paragraphs after the first one by manipulating the XML
        p_element = first_para._p
        parent = p_element.getparent()
        for p in list(parent.iterchildren()):
            if p != p_element and p.tag.endswith("}p"):
                parent.remove(p)


def apply_paragraph_formatting(
    paragraph: Any,
    spec: ParagraphSpec,
    default_font_size: Optional[float] = None,
) -> None:
    """Apply formatting from specification to a paragraph."""
    # Apply alignment
    alignment_map = {
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
        "JUSTIFY": PP_ALIGN.JUSTIFY,
        "LEFT": PP_ALIGN.LEFT,
    }
    alignment = spec.get("alignment")
    if alignment and alignment in alignment_map:
        paragraph.alignment = alignment_map[alignment]

    # Apply bullet
    if spec.get("bullet"):
        paragraph.level = spec.get("level", 0)
        # Enable bullet character (requires XML manipulation)
        pPr = paragraph._p.get_or_add_pPr()
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        # Remove any existing bullet definitions
        for elem in list(pPr):
            if elem.tag.startswith(ns + "bu"):
                pPr.remove(elem)
        # Add character bullet
        from lxml import etree
        buChar = etree.SubElement(pPr, f"{ns}buChar")
        buChar.set("char", "•")

    # Apply spacing
    space_before = spec.get("space_before")
    if space_before is not None:
        paragraph.space_before = Pt(space_before)

    space_after = spec.get("space_after")
    if space_after is not None:
        paragraph.space_after = Pt(space_after)

    # Apply line spacing
    line_spacing = spec.get("line_spacing")
    if line_spacing is not None:
        paragraph.line_spacing = Pt(line_spacing)

    # Set text content
    text = spec.get("text", "")
    run = paragraph.add_run()
    run.text = str(text)

    # Apply font formatting to run
    font = run.font
    font_name = spec.get("font_name")
    if font_name:
        font.name = font_name

    font_size = spec.get("font_size")
    if font_size is not None:
        font.size = Pt(font_size)
    elif default_font_size is not None:
        font.size = Pt(default_font_size)

    if spec.get("bold") is not None:
        font.bold = spec["bold"]

    if spec.get("italic") is not None:
        font.italic = spec["italic"]

    if spec.get("underline") is not None:
        font.underline = spec["underline"]

    # Apply color
    color = spec.get("color")
    if color:
        from pptx.dml.color import RGBColor
        try:
            font.color.rgb = RGBColor.from_string(color)
        except (ValueError, AttributeError):
            pass


def apply_replacement_to_shape(
    shape: BaseShape,
    paragraphs: List[ParagraphSpec],
    default_font_size: Optional[float] = None,
) -> None:
    """Apply replacement paragraphs to a shape."""
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return

    tf = shape.text_frame

    # Clear existing content
    clear_text_frame(shape)

    # Apply each paragraph
    for i, para_spec in enumerate(paragraphs):
        if i == 0:
            # Use existing first paragraph
            paragraph = tf.paragraphs[0]
        else:
            # Add new paragraph by adding a <a:p> element to the text frame
            from lxml import etree
            txBody = tf._txBody
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            # Create new paragraph element
            new_p = etree.SubElement(txBody, f"{ns}p")
            # Get the newly created paragraph through python-pptx
            paragraph = tf.paragraphs[i]

        apply_paragraph_formatting(paragraph, para_spec, default_font_size)


def get_shapes_for_slide(
    slide: Any,
    slide_idx: int,
) -> Tuple[Dict[str, BaseShape], Dict[str, float]]:
    """Get sorted shapes for a slide and their default font sizes."""
    shapes_with_positions = []
    for shape in slide.shapes:
        shapes_with_positions.extend(collect_shapes_with_absolute_positions(shape))

    if not shapes_with_positions:
        return {}, {}

    shape_data_list = [
        ShapeData(swp.shape, swp.absolute_left, swp.absolute_top, slide)
        for swp in shapes_with_positions
    ]

    sorted_shapes = sort_shapes_by_position(shape_data_list)

    shape_map: Dict[str, BaseShape] = {}
    font_sizes: Dict[str, float] = {}

    for idx, shape_data in enumerate(sorted_shapes):
        shape_id = f"shape-{idx}"
        shape_map[shape_id] = shape_data.shape
        if shape_data.default_font_size:
            font_sizes[shape_id] = shape_data.default_font_size

    return shape_map, font_sizes


def validate_replacements(
    replacements: ReplacementSpec,
    prs: Presentation,
) -> List[str]:
    """Validate replacement specification against presentation structure.

    Returns list of error messages (empty if valid).
    """
    errors = []

    for slide_id in replacements:
        if not validate_slide_id(slide_id):
            errors.append(f"Invalid slide ID format: {slide_id}")
            continue

        slide_idx = parse_slide_index(slide_id)
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            errors.append(f"Slide index out of range: {slide_id} (presentation has {len(prs.slides)} slides)")
            continue

        slide = prs.slides[slide_idx]
        shape_map, _ = get_shapes_for_slide(slide, slide_idx)

        slide_shapes = replacements[slide_id]
        if not isinstance(slide_shapes, dict):
            errors.append(f"{slide_id}: Expected dictionary of shape replacements")
            continue

        for shape_id, para_specs in slide_shapes.items():
            if not validate_shape_id(shape_id):
                errors.append(f"{slide_id}/{shape_id}: Invalid shape ID format")
                continue

            if shape_id not in shape_map:
                errors.append(f"{slide_id}/{shape_id}: Shape not found (available: {list(shape_map.keys())})")
                continue

            if not isinstance(para_specs, list):
                errors.append(f"{slide_id}/{shape_id}: Expected list of paragraph specifications")
                continue

            for i, para_spec in enumerate(para_specs):
                if not isinstance(para_spec, dict):
                    errors.append(f"{slide_id}/{shape_id}/paragraph-{i}: Expected dictionary")
                    continue

                if "text" not in para_spec:
                    errors.append(f"{slide_id}/{shape_id}/paragraph-{i}: Missing required 'text' field")

    return errors


def apply_replacements(
    pptx_path: str,
    json_path: str,
    output_path: str,
    clear_unspecified: bool = True,
) -> None:
    """Apply text replacements from JSON specification to a PowerPoint file.

    Args:
        pptx_path: Path to input PowerPoint file
        json_path: Path to JSON file with replacement specifications
        output_path: Path for output PowerPoint file
        clear_unspecified: If True, clear shapes not specified in JSON

    Raises:
        ValueError: If validation fails
        FileNotFoundError: If input files don't exist
    """
    pptx_path = Path(pptx_path)
    json_path = Path(json_path)
    output_path = Path(output_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

    if not json_path.exists():
        raise FileNotFoundError(f"JSON file not found: {json_path}")

    # Load JSON replacements
    with open(json_path, "r", encoding="utf-8") as f:
        replacements: ReplacementSpec = json.load(f)

    # Load presentation
    prs = Presentation(str(pptx_path))

    # Validate replacements
    errors = validate_replacements(replacements, prs)
    if errors:
        raise ValueError("Validation failed:\n" + "\n".join(f"  - {e}" for e in errors))

    # Apply replacements slide by slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_id = f"slide-{slide_idx}"
        shape_map, font_sizes = get_shapes_for_slide(slide, slide_idx)

        slide_replacements = replacements.get(slide_id, {})

        for shape_id, shape in shape_map.items():
            if shape_id in slide_replacements:
                # Apply specified replacement
                para_specs = slide_replacements[shape_id]
                default_font_size = font_sizes.get(shape_id)
                apply_replacement_to_shape(shape, para_specs, default_font_size)
            elif clear_unspecified:
                # Clear unspecified shapes
                clear_text_frame(shape)

    # Save output
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
