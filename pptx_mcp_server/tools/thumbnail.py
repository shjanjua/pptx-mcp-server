"""
Create thumbnail grids from PowerPoint presentations.

This module provides functionality to:
- Render PowerPoint slides to images
- Create thumbnail grid images for visual inspection
- Optionally outline text placeholders
- Support for large presentations with multiple grid pages
"""

import os
import subprocess
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image, ImageDraw


# Constants
MAX_SLIDES_PER_GRID = 60  # Maximum slides per grid image
DPI = 150  # Resolution for rendering
GRID_WIDTH = 2400  # Target grid width in pixels


def get_slide_dimensions(pptx_path: Path) -> Tuple[int, int]:
    """Get slide dimensions from a PowerPoint presentation.

    Returns:
        Tuple of (width, height) in pixels at specified DPI
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(pptx_path))
    width_inches = prs.slide_width / Inches(1)
    height_inches = prs.slide_height / Inches(1)

    width_px = int(width_inches * DPI)
    height_px = int(height_inches * DPI)

    return width_px, height_px


def render_slides_to_images(
    pptx_path: Path,
    output_dir: Path,
    dpi: int = DPI,
) -> List[Path]:
    """Render all slides to PNG images using LibreOffice.

    Args:
        pptx_path: Path to the PowerPoint file
        output_dir: Directory to save rendered images
        dpi: Resolution for rendering

    Returns:
        List of paths to rendered PNG images
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    # Use LibreOffice to convert to PDF first, then to images
    # This is more reliable than direct image export
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        # Convert to PDF
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(temp_path),
                str(pptx_path),
            ],
            capture_output=True,
            timeout=120,
        )

        if result.returncode != 0:
            raise RuntimeError(f"Failed to convert to PDF: {result.stderr.decode()}")

        pdf_path = temp_path / f"{pptx_path.stem}.pdf"
        if not pdf_path.exists():
            raise RuntimeError("PDF conversion failed - output file not created")

        # Convert PDF to images using ImageMagick or pdftoppm
        try:
            # Try pdftoppm (poppler-utils) first - faster and better quality
            result = subprocess.run(
                [
                    "pdftoppm",
                    "-png",
                    "-r", str(dpi),
                    str(pdf_path),
                    str(output_dir / "slide"),
                ],
                capture_output=True,
                timeout=300,
            )

            if result.returncode == 0:
                # pdftoppm names files as slide-01.png, slide-02.png, etc.
                images = sorted(output_dir.glob("slide-*.png"))
                if images:
                    return list(images)

        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Fallback to ImageMagick convert
        try:
            result = subprocess.run(
                [
                    "convert",
                    "-density", str(dpi),
                    str(pdf_path),
                    str(output_dir / "slide-%02d.png"),
                ],
                capture_output=True,
                timeout=300,
            )

            if result.returncode == 0:
                images = sorted(output_dir.glob("slide-*.png"))
                if images:
                    return list(images)

        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        raise RuntimeError(
            "Could not render slides to images. "
            "Please install pdftoppm (poppler-utils) or ImageMagick."
        )


def get_text_placeholder_bounds(pptx_path: Path) -> List[List[Tuple[float, float, float, float]]]:
    """Get bounding boxes of text placeholders for each slide.

    Returns:
        List of lists, where each inner list contains (left, top, width, height)
        tuples for text shapes on that slide, in inches.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(pptx_path))
    all_bounds = []

    for slide in prs.slides:
        slide_bounds = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                if shape.text_frame.text.strip():
                    # Convert EMU to inches
                    left = shape.left / Inches(1)
                    top = shape.top / Inches(1)
                    width = shape.width / Inches(1)
                    height = shape.height / Inches(1)
                    slide_bounds.append((left, top, width, height))
        all_bounds.append(slide_bounds)

    return all_bounds


def create_grid_image(
    images: List[Path],
    cols: int,
    thumb_width: int,
    thumb_height: int,
    output_path: Path,
    start_index: int = 0,
    placeholder_bounds: Optional[List[List[Tuple[float, float, float, float]]]] = None,
    slide_dims: Optional[Tuple[int, int]] = None,
) -> None:
    """Create a grid image from slide thumbnails.

    Args:
        images: List of paths to slide images
        cols: Number of columns in the grid
        thumb_width: Width of each thumbnail
        thumb_height: Height of each thumbnail
        output_path: Path to save the grid image
        start_index: Starting slide index for numbering
        placeholder_bounds: Optional bounds for text placeholders to outline
        slide_dims: Original slide dimensions in pixels
    """
    num_images = len(images)
    rows = (num_images + cols - 1) // cols

    # Add padding
    padding = 10
    label_height = 20

    grid_width = cols * thumb_width + (cols + 1) * padding
    grid_height = rows * (thumb_height + label_height) + (rows + 1) * padding

    # Create grid image with white background
    grid = Image.new("RGB", (grid_width, grid_height), "white")
    draw = ImageDraw.Draw(grid)

    for i, image_path in enumerate(images):
        row = i // cols
        col = i % cols

        x = col * thumb_width + (col + 1) * padding
        y = row * (thumb_height + label_height) + (row + 1) * padding

        # Load and resize thumbnail
        try:
            thumb = Image.open(image_path)
            thumb = thumb.resize((thumb_width, thumb_height), Image.Resampling.LANCZOS)

            # Draw placeholder outlines if requested
            if placeholder_bounds and slide_dims:
                slide_idx = start_index + i
                if slide_idx < len(placeholder_bounds):
                    thumb_draw = ImageDraw.Draw(thumb)
                    orig_width, orig_height = slide_dims
                    scale_x = thumb_width / orig_width
                    scale_y = thumb_height / orig_height

                    for bounds in placeholder_bounds[slide_idx]:
                        left, top, width, height = bounds
                        # Convert inches to pixels
                        px_left = int(left * DPI * scale_x)
                        px_top = int(top * DPI * scale_y)
                        px_right = int((left + width) * DPI * scale_x)
                        px_bottom = int((top + height) * DPI * scale_y)

                        thumb_draw.rectangle(
                            [px_left, px_top, px_right, px_bottom],
                            outline="red",
                            width=2,
                        )

            # Paste thumbnail
            grid.paste(thumb, (x, y))

            # Draw slide number label
            slide_num = start_index + i
            label = f"Slide {slide_num}"
            draw.text(
                (x + thumb_width // 2, y + thumb_height + 2),
                label,
                fill="black",
                anchor="ma",
            )

        except Exception as e:
            # Draw error placeholder
            draw.rectangle([x, y, x + thumb_width, y + thumb_height], fill="lightgray")
            draw.text(
                (x + thumb_width // 2, y + thumb_height // 2),
                f"Error: {str(e)[:20]}",
                fill="red",
                anchor="mm",
            )

    # Save grid
    output_path.parent.mkdir(parents=True, exist_ok=True)
    grid.save(str(output_path), "PNG")


def create_thumbnail_grids(
    pptx_path: Path,
    output_prefix: str = "thumbnails",
    cols: int = 5,
    outline_placeholders: bool = False,
) -> List[str]:
    """Create thumbnail grid images from a PowerPoint presentation.

    Args:
        pptx_path: Path to the PowerPoint file
        output_prefix: Prefix for output files
        cols: Number of columns in the grid (3-6)
        outline_placeholders: Whether to outline text placeholders

    Returns:
        List of paths to created grid images
    """
    pptx_path = Path(pptx_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

    # Clamp cols to valid range
    cols = min(max(cols, 3), 6)

    # Get slide dimensions
    slide_width, slide_height = get_slide_dimensions(pptx_path)

    # Calculate thumbnail dimensions
    thumb_width = GRID_WIDTH // (cols + 1)
    thumb_height = int(thumb_width * slide_height / slide_width)

    # Render slides to images
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        slide_images = render_slides_to_images(pptx_path, temp_path)

        if not slide_images:
            raise RuntimeError("No slide images were rendered")

        # Get placeholder bounds if needed
        placeholder_bounds = None
        if outline_placeholders:
            placeholder_bounds = get_text_placeholder_bounds(pptx_path)

        # Create grids
        grid_files = []
        num_slides = len(slide_images)
        num_grids = (num_slides + MAX_SLIDES_PER_GRID - 1) // MAX_SLIDES_PER_GRID

        for grid_idx in range(num_grids):
            start_idx = grid_idx * MAX_SLIDES_PER_GRID
            end_idx = min(start_idx + MAX_SLIDES_PER_GRID, num_slides)

            grid_images = slide_images[start_idx:end_idx]

            if num_grids > 1:
                output_path = Path(f"{output_prefix}-{grid_idx + 1}.png")
            else:
                output_path = Path(f"{output_prefix}.png")

            create_grid_image(
                grid_images,
                cols,
                thumb_width,
                thumb_height,
                output_path,
                start_index=start_idx,
                placeholder_bounds=placeholder_bounds,
                slide_dims=(slide_width, slide_height),
            )

            grid_files.append(str(output_path))

        return grid_files
