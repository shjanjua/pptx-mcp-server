"""
Create new PowerPoint presentations from scratch.

This module provides functionality to:
- Create new presentations with specified dimensions
- Add slides with text boxes, shapes, and images
- Apply formatting (fonts, colors, alignment, bullets)
- Support for common layouts (title slide, content slide, etc.)
"""

import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# Type aliases
SlideSpec = Dict[str, Any]
ShapeSpec = Dict[str, Any]
PresentationSpec = Dict[str, Any]


# Standard layouts
LAYOUTS = {
    "16:9": {"width": 13.333, "height": 7.5},
    "4:3": {"width": 10.0, "height": 7.5},
    "widescreen": {"width": 13.333, "height": 7.5},
    "standard": {"width": 10.0, "height": 7.5},
}

# Alignment mapping
ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

VALIGN_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def create_presentation_from_spec(spec: PresentationSpec, output_path: Path) -> str:
    """Create a new PowerPoint presentation from a JSON specification.

    Args:
        spec: Presentation specification dictionary
        output_path: Path to save the presentation

    Returns:
        Success message with details
    """
    # Create presentation
    prs = Presentation()

    # Set dimensions
    layout = spec.get("layout", "16:9")
    if layout in LAYOUTS:
        dims = LAYOUTS[layout]
        prs.slide_width = Inches(dims["width"])
        prs.slide_height = Inches(dims["height"])
    elif "width" in spec and "height" in spec:
        prs.slide_width = Inches(spec["width"])
        prs.slide_height = Inches(spec["height"])

    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Process slides
    slides_spec = spec.get("slides", [])
    for slide_spec in slides_spec:
        slide = prs.slides.add_slide(blank_layout)
        _add_slide_content(slide, slide_spec, prs)

    # Save
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    return f"Created presentation with {len(slides_spec)} slides: {output_path}"


def _add_slide_content(slide: Any, spec: SlideSpec, prs: Presentation) -> None:
    """Add content to a slide based on specification."""
    # Set background color if specified
    background = spec.get("background")
    if background:
        if isinstance(background, str) and background.startswith("#"):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(background.lstrip("#"))
        elif isinstance(background, dict) and "color" in background:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(background["color"].lstrip("#"))

    # Add shapes
    shapes = spec.get("shapes", [])
    for shape_spec in shapes:
        _add_shape(slide, shape_spec, prs)


def _add_shape(slide: Any, spec: ShapeSpec, prs: Presentation) -> None:
    """Add a shape to the slide based on specification."""
    shape_type = spec.get("type", "textbox")

    # Get position and size (in inches)
    left = Inches(spec.get("left", 0.5))
    top = Inches(spec.get("top", 0.5))
    width = Inches(spec.get("width", 5))
    height = Inches(spec.get("height", 1))

    if shape_type == "textbox":
        shape = slide.shapes.add_textbox(left, top, width, height)
        _apply_text_content(shape, spec)

    elif shape_type == "rectangle":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        _apply_shape_fill(shape, spec)
        if "text" in spec or "paragraphs" in spec:
            _apply_text_content(shape, spec)

    elif shape_type == "rounded_rectangle":
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        _apply_shape_fill(shape, spec)
        if "text" in spec or "paragraphs" in spec:
            _apply_text_content(shape, spec)

    elif shape_type == "oval":
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        _apply_shape_fill(shape, spec)
        if "text" in spec or "paragraphs" in spec:
            _apply_text_content(shape, spec)

    elif shape_type == "image":
        image_path = spec.get("path") or spec.get("src")
        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(image_path, left, top, width, height)

    elif shape_type == "line":
        # For lines, width/height represent the end point offset
        end_left = left + width
        end_top = top + height
        shape = slide.shapes.add_connector(
            1,  # Straight connector
            left, top,
            end_left, end_top
        )
        if "color" in spec:
            shape.line.color.rgb = RGBColor.from_string(spec["color"].lstrip("#"))
        if "line_width" in spec:
            shape.line.width = Pt(spec["line_width"])


def _apply_shape_fill(shape: Any, spec: ShapeSpec) -> None:
    """Apply fill properties to a shape."""
    fill_color = spec.get("fill") or spec.get("fill_color") or spec.get("background")
    if fill_color:
        shape.fill.solid()
        if isinstance(fill_color, str):
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))

    # Border/line
    border = spec.get("border") or spec.get("line")
    if border:
        if isinstance(border, dict):
            if "color" in border:
                shape.line.color.rgb = RGBColor.from_string(border["color"].lstrip("#"))
            if "width" in border:
                shape.line.width = Pt(border["width"])
        elif border is False or border == "none":
            shape.line.fill.background()
    elif spec.get("no_border") or spec.get("no_line"):
        shape.line.fill.background()


def _apply_text_content(shape: Any, spec: ShapeSpec) -> None:
    """Apply text content and formatting to a shape."""
    tf = shape.text_frame
    tf.word_wrap = spec.get("word_wrap", True)

    # Vertical alignment
    valign = spec.get("valign", "top")
    if valign in VALIGN_MAP:
        tf.anchor = VALIGN_MAP[valign]

    # Margins
    if "margin" in spec:
        margin = spec["margin"]
        if isinstance(margin, (int, float)):
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
        elif isinstance(margin, list) and len(margin) == 4:
            tf.margin_top = Inches(margin[0])
            tf.margin_right = Inches(margin[1])
            tf.margin_bottom = Inches(margin[2])
            tf.margin_left = Inches(margin[3])

    # Handle paragraphs
    paragraphs = spec.get("paragraphs", [])
    if not paragraphs and "text" in spec:
        # Simple text - convert to single paragraph
        paragraphs = [{"text": spec["text"]}]
        # Inherit formatting from shape spec
        for key in ["font_size", "font_name", "bold", "italic", "color", "alignment"]:
            if key in spec:
                paragraphs[0][key] = spec[key]

    for i, para_spec in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        _apply_paragraph(p, para_spec)


def _apply_paragraph(paragraph: Any, spec: Dict[str, Any]) -> None:
    """Apply paragraph formatting and content."""
    # Alignment
    align = spec.get("alignment", spec.get("align", "left"))
    if align in ALIGN_MAP:
        paragraph.alignment = ALIGN_MAP[align]

    # Spacing
    if "space_before" in spec:
        paragraph.space_before = Pt(spec["space_before"])
    if "space_after" in spec:
        paragraph.space_after = Pt(spec["space_after"])
    if "line_spacing" in spec:
        paragraph.line_spacing = Pt(spec["line_spacing"])

    # Bullet
    if spec.get("bullet"):
        paragraph.level = spec.get("level", 0)
        # Enable bullet
        pPr = paragraph._p.get_or_add_pPr()
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        # Remove existing bullet definitions
        for elem in list(pPr):
            if elem.tag.startswith(ns + "bu"):
                pPr.remove(elem)
        # Add bullet character
        from lxml import etree
        buChar = etree.SubElement(pPr, f"{ns}buChar")
        buChar.set("char", spec.get("bullet_char", "•"))

    # Text content
    text = spec.get("text", "")
    run = paragraph.add_run()
    run.text = str(text)

    # Font formatting
    font = run.font
    if "font_name" in spec:
        font.name = spec["font_name"]
    if "font_size" in spec:
        font.size = Pt(spec["font_size"])
    if "bold" in spec:
        font.bold = spec["bold"]
    if "italic" in spec:
        font.italic = spec["italic"]
    if "underline" in spec:
        font.underline = spec["underline"]
    if "color" in spec:
        color = spec["color"]
        if isinstance(color, str):
            font.color.rgb = RGBColor.from_string(color.lstrip("#"))


def create_presentation(
    output_path: str,
    spec: Optional[Union[str, Dict]] = None,
    layout: str = "16:9",
    slides: Optional[List[Dict]] = None,
) -> str:
    """Create a new PowerPoint presentation.

    Can be called with either:
    1. A full spec dict/JSON string
    2. Individual parameters (layout + slides)

    Args:
        output_path: Path to save the presentation
        spec: Full presentation specification (dict or JSON string)
        layout: Slide layout ("16:9", "4:3", etc.)
        slides: List of slide specifications

    Returns:
        Success message
    """
    if spec:
        if isinstance(spec, str):
            spec = json.loads(spec)
    else:
        spec = {
            "layout": layout,
            "slides": slides or [],
        }

    return create_presentation_from_spec(spec, Path(output_path))
