"""
Extract structured text content from PowerPoint presentations.

This module provides functionality to:
- Extract all text content from PowerPoint shapes
- Preserve paragraph formatting (alignment, bullets, fonts, spacing)
- Handle nested GroupShapes recursively with correct absolute positions
- Sort shapes by visual position on slides
- Filter out slide numbers and non-content placeholders
- Export to JSON with clean, structured data
"""

import json
import platform
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape

# Type aliases
JsonValue = Union[str, int, float, bool, None]
ParagraphDict = Dict[str, JsonValue]
ShapeDict = Dict[str, Union[str, float, bool, List[ParagraphDict], List[str], Dict[str, Any], None]]
InventoryData = Dict[str, Dict[str, "ShapeData"]]
InventoryDict = Dict[str, Dict[str, ShapeDict]]


@dataclass
class ShapeWithPosition:
    """A shape with its absolute position on the slide."""
    shape: BaseShape
    absolute_left: int  # in EMUs
    absolute_top: int  # in EMUs


class ParagraphData:
    """Data structure for paragraph properties extracted from a PowerPoint paragraph."""

    def __init__(self, paragraph: Any):
        """Initialize from a PowerPoint paragraph object."""
        self.text: str = paragraph.text.strip()
        self.bullet: bool = False
        self.level: Optional[int] = None
        self.alignment: Optional[str] = None
        self.space_before: Optional[float] = None
        self.space_after: Optional[float] = None
        self.font_name: Optional[str] = None
        self.font_size: Optional[float] = None
        self.bold: Optional[bool] = None
        self.italic: Optional[bool] = None
        self.underline: Optional[bool] = None
        self.color: Optional[str] = None
        self.theme_color: Optional[str] = None
        self.line_spacing: Optional[float] = None

        # Check for bullet formatting
        if hasattr(paragraph, "_p") and paragraph._p is not None and paragraph._p.pPr is not None:
            pPr = paragraph._p.pPr
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            if pPr.find(f"{ns}buChar") is not None or pPr.find(f"{ns}buAutoNum") is not None:
                self.bullet = True
                if hasattr(paragraph, "level"):
                    self.level = paragraph.level

        # Add alignment if not LEFT (default)
        if hasattr(paragraph, "alignment") and paragraph.alignment is not None:
            alignment_map = {
                PP_ALIGN.CENTER: "CENTER",
                PP_ALIGN.RIGHT: "RIGHT",
                PP_ALIGN.JUSTIFY: "JUSTIFY",
            }
            if paragraph.alignment in alignment_map:
                self.alignment = alignment_map[paragraph.alignment]

        # Add spacing properties if set
        if hasattr(paragraph, "space_before") and paragraph.space_before:
            self.space_before = paragraph.space_before.pt
        if hasattr(paragraph, "space_after") and paragraph.space_after:
            self.space_after = paragraph.space_after.pt

        # Extract font properties from first run
        if paragraph.runs:
            first_run = paragraph.runs[0]
            if hasattr(first_run, "font"):
                font = first_run.font
                if font.name:
                    self.font_name = font.name
                if font.size:
                    self.font_size = font.size.pt
                if font.bold is not None:
                    self.bold = font.bold
                if font.italic is not None:
                    self.italic = font.italic
                if font.underline is not None:
                    self.underline = font.underline

                # Handle color
                try:
                    if font.color.rgb:
                        self.color = str(font.color.rgb)
                except (AttributeError, TypeError):
                    try:
                        if font.color.theme_color:
                            self.theme_color = font.color.theme_color.name
                    except (AttributeError, TypeError):
                        pass

        # Add line spacing if set
        if hasattr(paragraph, "line_spacing") and paragraph.line_spacing is not None:
            if hasattr(paragraph.line_spacing, "pt"):
                self.line_spacing = round(paragraph.line_spacing.pt, 2)
            else:
                font_size = self.font_size if self.font_size else 12.0
                self.line_spacing = round(paragraph.line_spacing * font_size, 2)

    def to_dict(self) -> ParagraphDict:
        """Convert to dictionary for JSON serialization, excluding None values."""
        result: ParagraphDict = {"text": self.text}

        if self.bullet:
            result["bullet"] = self.bullet
        if self.level is not None:
            result["level"] = self.level
        if self.alignment:
            result["alignment"] = self.alignment
        if self.space_before is not None:
            result["space_before"] = self.space_before
        if self.space_after is not None:
            result["space_after"] = self.space_after
        if self.font_name:
            result["font_name"] = self.font_name
        if self.font_size is not None:
            result["font_size"] = self.font_size
        if self.bold is not None:
            result["bold"] = self.bold
        if self.italic is not None:
            result["italic"] = self.italic
        if self.underline is not None:
            result["underline"] = self.underline
        if self.color:
            result["color"] = self.color
        if self.theme_color:
            result["theme_color"] = self.theme_color
        if self.line_spacing is not None:
            result["line_spacing"] = self.line_spacing

        return result


class ShapeData:
    """Data structure for shape properties extracted from a PowerPoint shape."""

    @staticmethod
    def emu_to_inches(emu: int) -> float:
        """Convert EMUs (English Metric Units) to inches."""
        return emu / 914400.0

    @staticmethod
    def inches_to_pixels(inches: float, dpi: int = 96) -> int:
        """Convert inches to pixels at given DPI."""
        return int(inches * dpi)

    @staticmethod
    def get_font_path(font_name: str) -> Optional[str]:
        """Get the font file path for a given font name."""
        system = platform.system()
        font_variations = [font_name, font_name.lower(), font_name.replace(" ", ""), font_name.replace(" ", "-")]

        if system == "Darwin":  # macOS
            font_dirs = ["/System/Library/Fonts/", "/Library/Fonts/", "~/Library/Fonts/"]
            extensions = [".ttf", ".otf", ".ttc", ".dfont"]
        else:  # Linux
            font_dirs = ["/usr/share/fonts/truetype/", "/usr/local/share/fonts/", "~/.fonts/"]
            extensions = [".ttf", ".otf"]

        for font_dir in font_dirs:
            font_dir_path = Path(font_dir).expanduser()
            if not font_dir_path.exists():
                continue

            for variant in font_variations:
                for ext in extensions:
                    font_path = font_dir_path / f"{variant}{ext}"
                    if font_path.exists():
                        return str(font_path)

            try:
                for file_path in font_dir_path.iterdir():
                    if file_path.is_file():
                        file_name_lower = file_path.name.lower()
                        font_name_lower = font_name.lower().replace(" ", "")
                        if font_name_lower in file_name_lower and any(
                            file_name_lower.endswith(ext) for ext in extensions
                        ):
                            return str(file_path)
            except (OSError, PermissionError):
                continue

        return None

    @staticmethod
    def get_slide_dimensions(slide: Any) -> Tuple[Optional[int], Optional[int]]:
        """Get slide dimensions from slide object."""
        try:
            prs = slide.part.package.presentation_part.presentation
            return prs.slide_width, prs.slide_height
        except (AttributeError, TypeError):
            return None, None

    @staticmethod
    def get_default_font_size(shape: BaseShape, slide_layout: Any) -> Optional[float]:
        """Extract default font size from slide layout for a placeholder shape."""
        try:
            if not hasattr(shape, "placeholder_format"):
                return None

            shape_type = shape.placeholder_format.type
            for layout_placeholder in slide_layout.placeholders:
                if layout_placeholder.placeholder_format.type == shape_type:
                    for elem in layout_placeholder.element.iter():
                        if "defRPr" in elem.tag and (sz := elem.get("sz")):
                            return float(sz) / 100.0
                    break
        except Exception:
            pass
        return None

    def __init__(
        self,
        shape: BaseShape,
        absolute_left: Optional[int] = None,
        absolute_top: Optional[int] = None,
        slide: Optional[Any] = None,
    ):
        """Initialize from a PowerPoint shape object."""
        self.shape = shape
        self.shape_id: str = ""

        self.slide_width_emu, self.slide_height_emu = (
            self.get_slide_dimensions(slide) if slide else (None, None)
        )

        self.placeholder_type: Optional[str] = None
        self.default_font_size: Optional[float] = None
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            if shape.placeholder_format and shape.placeholder_format.type:
                self.placeholder_type = str(shape.placeholder_format.type).split(".")[-1].split(" ")[0]
                if slide and hasattr(slide, "slide_layout"):
                    self.default_font_size = self.get_default_font_size(shape, slide.slide_layout)

        left_emu = absolute_left if absolute_left is not None else (shape.left if hasattr(shape, "left") else 0)
        top_emu = absolute_top if absolute_top is not None else (shape.top if hasattr(shape, "top") else 0)

        self.left: float = round(self.emu_to_inches(left_emu), 2)
        self.top: float = round(self.emu_to_inches(top_emu), 2)
        self.width: float = round(self.emu_to_inches(shape.width if hasattr(shape, "width") else 0), 2)
        self.height: float = round(self.emu_to_inches(shape.height if hasattr(shape, "height") else 0), 2)

        self.left_emu = left_emu
        self.top_emu = top_emu
        self.width_emu = shape.width if hasattr(shape, "width") else 0
        self.height_emu = shape.height if hasattr(shape, "height") else 0

        self.frame_overflow_bottom: Optional[float] = None
        self.slide_overflow_right: Optional[float] = None
        self.slide_overflow_bottom: Optional[float] = None
        self.overlapping_shapes: Dict[str, float] = {}
        self.warnings: List[str] = []

        self._estimate_frame_overflow()
        self._calculate_slide_overflow()
        self._detect_bullet_issues()

    @property
    def paragraphs(self) -> List[ParagraphData]:
        """Calculate paragraphs from the shape's text frame."""
        if not self.shape or not hasattr(self.shape, "text_frame"):
            return []

        paragraphs = []
        for paragraph in self.shape.text_frame.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(ParagraphData(paragraph))
        return paragraphs

    def _get_default_font_size(self) -> int:
        """Get default font size from theme text styles or use conservative default."""
        try:
            if not (hasattr(self.shape, "part") and hasattr(self.shape.part, "slide_layout")):
                return 14

            slide_master = self.shape.part.slide_layout.slide_master
            if not hasattr(slide_master, "element"):
                return 14

            style_name = "bodyStyle"
            if self.placeholder_type and "TITLE" in self.placeholder_type:
                style_name = "titleStyle"

            for child in slide_master.element.iter():
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == style_name:
                    for elem in child.iter():
                        if "sz" in elem.attrib:
                            return int(elem.attrib["sz"]) // 100
        except Exception:
            pass

        return 14

    def _get_usable_dimensions(self, text_frame) -> Tuple[int, int]:
        """Get usable width and height in pixels after accounting for margins."""
        margins = {"top": 0.05, "bottom": 0.05, "left": 0.1, "right": 0.1}

        if hasattr(text_frame, "margin_top") and text_frame.margin_top:
            margins["top"] = self.emu_to_inches(text_frame.margin_top)
        if hasattr(text_frame, "margin_bottom") and text_frame.margin_bottom:
            margins["bottom"] = self.emu_to_inches(text_frame.margin_bottom)
        if hasattr(text_frame, "margin_left") and text_frame.margin_left:
            margins["left"] = self.emu_to_inches(text_frame.margin_left)
        if hasattr(text_frame, "margin_right") and text_frame.margin_right:
            margins["right"] = self.emu_to_inches(text_frame.margin_right)

        usable_width = self.width - margins["left"] - margins["right"]
        usable_height = self.height - margins["top"] - margins["bottom"]

        return self.inches_to_pixels(usable_width), self.inches_to_pixels(usable_height)

    def _wrap_text_line(self, line: str, max_width_px: int, draw, font) -> List[str]:
        """Wrap a single line of text to fit within max_width_px."""
        if not line:
            return [""]

        if draw.textlength(line, font=font) <= max_width_px:
            return [line]

        wrapped = []
        words = line.split(" ")
        current_line = ""

        for word in words:
            test_line = current_line + (" " if current_line else "") + word
            if draw.textlength(test_line, font=font) <= max_width_px:
                current_line = test_line
            else:
                if current_line:
                    wrapped.append(current_line)
                current_line = word

        if current_line:
            wrapped.append(current_line)

        return wrapped

    def _estimate_frame_overflow(self) -> None:
        """Estimate if text overflows the shape bounds using PIL text measurement."""
        if not self.shape or not hasattr(self.shape, "text_frame"):
            return

        text_frame = self.shape.text_frame
        if not text_frame or not text_frame.paragraphs:
            return

        usable_width_px, usable_height_px = self._get_usable_dimensions(text_frame)
        if usable_width_px <= 0 or usable_height_px <= 0:
            return

        dummy_img = Image.new("RGB", (1, 1))
        draw = ImageDraw.Draw(dummy_img)

        default_font_size = self._get_default_font_size()
        total_height_px = 0

        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            if not paragraph.text.strip():
                continue

            para_data = ParagraphData(paragraph)
            font_name = para_data.font_name or "Arial"
            font_size = int(para_data.font_size or default_font_size)

            font = None
            font_path = self.get_font_path(font_name)
            if font_path:
                try:
                    font = ImageFont.truetype(font_path, size=font_size)
                except Exception:
                    font = ImageFont.load_default()
            else:
                font = ImageFont.load_default()

            all_wrapped_lines = []
            for line in paragraph.text.split("\n"):
                wrapped = self._wrap_text_line(line, usable_width_px, draw, font)
                all_wrapped_lines.extend(wrapped)

            if all_wrapped_lines:
                if para_data.line_spacing:
                    line_height_px = para_data.line_spacing * 96 / 72
                else:
                    line_height_px = font_size * 96 / 72

                if para_idx > 0 and para_data.space_before:
                    total_height_px += para_data.space_before * 96 / 72

                total_height_px += len(all_wrapped_lines) * line_height_px

                if para_data.space_after:
                    total_height_px += para_data.space_after * 96 / 72

        if total_height_px > usable_height_px:
            overflow_px = total_height_px - usable_height_px
            overflow_inches = round(overflow_px / 96.0, 2)
            if overflow_inches > 0.05:
                self.frame_overflow_bottom = overflow_inches

    def _calculate_slide_overflow(self) -> None:
        """Calculate if shape overflows the slide boundaries."""
        if self.slide_width_emu is None or self.slide_height_emu is None:
            return

        right_edge_emu = self.left_emu + self.width_emu
        if right_edge_emu > self.slide_width_emu:
            overflow_emu = right_edge_emu - self.slide_width_emu
            overflow_inches = round(self.emu_to_inches(overflow_emu), 2)
            if overflow_inches > 0.01:
                self.slide_overflow_right = overflow_inches

        bottom_edge_emu = self.top_emu + self.height_emu
        if bottom_edge_emu > self.slide_height_emu:
            overflow_emu = bottom_edge_emu - self.slide_height_emu
            overflow_inches = round(self.emu_to_inches(overflow_emu), 2)
            if overflow_inches > 0.01:
                self.slide_overflow_bottom = overflow_inches

    def _detect_bullet_issues(self) -> None:
        """Detect bullet point formatting issues in paragraphs."""
        if not self.shape or not hasattr(self.shape, "text_frame"):
            return

        text_frame = self.shape.text_frame
        if not text_frame or not text_frame.paragraphs:
            return

        bullet_symbols = ["•", "●", "○"]

        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if text and any(text.startswith(symbol + " ") for symbol in bullet_symbols):
                self.warnings.append("manual_bullet_symbol: use proper bullet formatting")
                break

    @property
    def has_any_issues(self) -> bool:
        """Check if shape has any issues (overflow, overlap, or warnings)."""
        return (
            self.frame_overflow_bottom is not None
            or self.slide_overflow_right is not None
            or self.slide_overflow_bottom is not None
            or len(self.overlapping_shapes) > 0
            or len(self.warnings) > 0
        )

    def to_dict(self) -> ShapeDict:
        """Convert to dictionary for JSON serialization."""
        result: ShapeDict = {
            "left": self.left,
            "top": self.top,
            "width": self.width,
            "height": self.height,
        }

        if self.placeholder_type:
            result["placeholder_type"] = self.placeholder_type

        if self.default_font_size:
            result["default_font_size"] = self.default_font_size

        overflow_data = {}
        if self.frame_overflow_bottom is not None:
            overflow_data["frame"] = {"overflow_bottom": self.frame_overflow_bottom}

        slide_overflow = {}
        if self.slide_overflow_right is not None:
            slide_overflow["overflow_right"] = self.slide_overflow_right
        if self.slide_overflow_bottom is not None:
            slide_overflow["overflow_bottom"] = self.slide_overflow_bottom
        if slide_overflow:
            overflow_data["slide"] = slide_overflow

        if overflow_data:
            result["overflow"] = overflow_data

        if self.overlapping_shapes:
            result["overlap"] = {"overlapping_shapes": self.overlapping_shapes}

        if self.warnings:
            result["warnings"] = self.warnings

        result["paragraphs"] = [para.to_dict() for para in self.paragraphs]

        return result


def is_valid_shape(shape: BaseShape) -> bool:
    """Check if a shape contains meaningful text content."""
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return False

    text = shape.text_frame.text.strip()
    if not text:
        return False

    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
        if shape.placeholder_format and shape.placeholder_format.type:
            placeholder_type = str(shape.placeholder_format.type).split(".")[-1].split(" ")[0]
            if placeholder_type == "SLIDE_NUMBER":
                return False
            if placeholder_type == "FOOTER" and text.isdigit():
                return False

    return True


def collect_shapes_with_absolute_positions(
    shape: BaseShape, parent_left: int = 0, parent_top: int = 0
) -> List[ShapeWithPosition]:
    """Recursively collect all shapes with valid text, calculating absolute positions."""
    if hasattr(shape, "shapes"):  # GroupShape
        result = []
        group_left = shape.left if hasattr(shape, "left") else 0
        group_top = shape.top if hasattr(shape, "top") else 0

        abs_group_left = parent_left + group_left
        abs_group_top = parent_top + group_top

        for child in shape.shapes:
            result.extend(collect_shapes_with_absolute_positions(child, abs_group_left, abs_group_top))
        return result

    if is_valid_shape(shape):
        shape_left = shape.left if hasattr(shape, "left") else 0
        shape_top = shape.top if hasattr(shape, "top") else 0

        return [ShapeWithPosition(shape=shape, absolute_left=parent_left + shape_left, absolute_top=parent_top + shape_top)]

    return []


def sort_shapes_by_position(shapes: List[ShapeData]) -> List[ShapeData]:
    """Sort shapes by visual position (top-to-bottom, left-to-right)."""
    if not shapes:
        return shapes

    shapes = sorted(shapes, key=lambda s: (s.top, s.left))

    result = []
    row = [shapes[0]]
    row_top = shapes[0].top

    for shape in shapes[1:]:
        if abs(shape.top - row_top) <= 0.5:
            row.append(shape)
        else:
            result.extend(sorted(row, key=lambda s: s.left))
            row = [shape]
            row_top = shape.top

    result.extend(sorted(row, key=lambda s: s.left))
    return result


def calculate_overlap(
    rect1: Tuple[float, float, float, float],
    rect2: Tuple[float, float, float, float],
    tolerance: float = 0.05,
) -> Tuple[bool, float]:
    """Calculate if and how much two rectangles overlap."""
    left1, top1, w1, h1 = rect1
    left2, top2, w2, h2 = rect2

    overlap_width = min(left1 + w1, left2 + w2) - max(left1, left2)
    overlap_height = min(top1 + h1, top2 + h2) - max(top1, top2)

    if overlap_width > tolerance and overlap_height > tolerance:
        overlap_area = overlap_width * overlap_height
        return True, round(overlap_area, 2)

    return False, 0


def detect_overlaps(shapes: List[ShapeData]) -> None:
    """Detect overlapping shapes and update their overlapping_shapes dictionaries."""
    n = len(shapes)

    for i in range(n):
        for j in range(i + 1, n):
            shape1 = shapes[i]
            shape2 = shapes[j]

            rect1 = (shape1.left, shape1.top, shape1.width, shape1.height)
            rect2 = (shape2.left, shape2.top, shape2.width, shape2.height)

            overlaps, overlap_area = calculate_overlap(rect1, rect2)

            if overlaps:
                shape1.overlapping_shapes[shape2.shape_id] = overlap_area
                shape2.overlapping_shapes[shape1.shape_id] = overlap_area


def extract_text_inventory(
    pptx_path: Path, prs: Optional[Any] = None, issues_only: bool = False
) -> InventoryData:
    """Extract text content from all slides in a PowerPoint presentation.

    Returns a nested dictionary: {slide-N: {shape-N: ShapeData}}
    """
    if prs is None:
        prs = Presentation(str(pptx_path))
    inventory: InventoryData = {}

    for slide_idx, slide in enumerate(prs.slides):
        shapes_with_positions = []
        for shape in slide.shapes:
            shapes_with_positions.extend(collect_shapes_with_absolute_positions(shape))

        if not shapes_with_positions:
            continue

        shape_data_list = [
            ShapeData(swp.shape, swp.absolute_left, swp.absolute_top, slide)
            for swp in shapes_with_positions
        ]

        sorted_shapes = sort_shapes_by_position(shape_data_list)
        for idx, shape_data in enumerate(sorted_shapes):
            shape_data.shape_id = f"shape-{idx}"

        if len(sorted_shapes) > 1:
            detect_overlaps(sorted_shapes)

        if issues_only:
            sorted_shapes = [sd for sd in sorted_shapes if sd.has_any_issues]

        if not sorted_shapes:
            continue

        inventory[f"slide-{slide_idx}"] = {shape_data.shape_id: shape_data for shape_data in sorted_shapes}

    return inventory


def get_inventory_as_dict(pptx_path: Path, issues_only: bool = False) -> InventoryDict:
    """Extract text inventory and return as JSON-serializable dictionaries."""
    inventory = extract_text_inventory(pptx_path, issues_only=issues_only)

    dict_inventory: InventoryDict = {}
    for slide_key, shapes in inventory.items():
        dict_inventory[slide_key] = {
            shape_key: shape_data.to_dict() for shape_key, shape_data in shapes.items()
        }

    return dict_inventory


def save_inventory(inventory: InventoryData, output_path: Path) -> None:
    """Save inventory to JSON file with proper formatting."""
    json_inventory: InventoryDict = {}
    for slide_key, shapes in inventory.items():
        json_inventory[slide_key] = {
            shape_key: shape_data.to_dict() for shape_key, shape_data in shapes.items()
        }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(json_inventory, f, indent=2, ensure_ascii=False)
